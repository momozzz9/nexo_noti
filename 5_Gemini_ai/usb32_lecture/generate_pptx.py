# -*- coding: utf-8 -*-
"""
USB 3.2 Rev 1.1 Technical Presentation Generator (Ultra-Premium Version)
삼성전자 HW 개발자 교육용 고품격 디자인 PPTX 자동 생성 스크립트
"""

import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: 'python-pptx' 라이브러리가 필요합니다. 'pip install python-pptx'를 실행하세요.")
    exit()

def create_premium_presentation():
    prs = Presentation()

    # 슬라이드 크기 설정 (Widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    COLOR_BG = RGBColor(10, 10, 15)       # 매우 어두운 남색
    COLOR_PRIMARY = RGBColor(59, 130, 246) # 삼성 블루 (Modern)
    COLOR_TEXT = RGBColor(248, 250, 252)   # 연한 회색/흰색
    COLOR_DIM = RGBColor(148, 163, 184)    # 어두운 회색
    COLOR_ACCENT = RGBColor(16, 185, 129)  # 성공/에메랄드

    def set_slide_background(slide):
        """슬라이드 배경을 어두운 색으로 설정"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_premium_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드
        set_slide_background(slide)
        
        # 하단 디자인 바
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.visible = False

        # 중앙 타이틀
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(60)
        p.font.name = '맑은 고딕' # Korean Windows Standard
        p.font.color.rgb = COLOR_TEXT

        # 서브타이틀
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.name = '맑은 고딕'
        p.font.color.rgb = COLOR_DIM

    def add_premium_content_slide(index, title, content_list, technical_note=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드
        set_slide_background(slide)

        # 상단 헤더 디자인
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.8))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = COLOR_PRIMARY
        header_bar.line.visible = False

        # 슬라이드 인덱스
        idx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(1), Inches(0.4))
        idx_box.text_frame.text = f"{index:02d}"
        idx_box.text_frame.paragraphs[0].font.size = Pt(14)
        idx_box.text_frame.paragraphs[0].font.bold = True
        idx_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(12), Inches(1))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = '맑은 고딕'
        p.font.color.rgb = COLOR_TEXT

        # 본문 내용 (좌측)
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(20)
            p.font.name = '맑은 고딕'
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(12)

        # 우측 기술 노트 박스 (Deep Dive)
        if technical_note:
            note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.2), Inches(4), Inches(4.5))
            note_box.fill.solid()
            note_box.fill.fore_color.rgb = RGBColor(25, 30, 45) # 약간 밝은 남색
            note_box.line.color.rgb = COLOR_PRIMARY
            note_box.line.width = Pt(1.5)

            note_tf = note_box.text_frame
            note_tf.margin_top = Inches(0.3)
            note_tf.margin_left = Inches(0.2)
            note_tf.margin_right = Inches(0.2)
            
            p = note_tf.add_paragraph()
            p.text = "TECHNICAL DEEP DIVE"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_ACCENT
            p.space_after = Pt(10)

            p = note_tf.add_paragraph()
            p.text = technical_note
            p.font.size = Pt(14)
            p.font.name = '맑은 고딕'
            p.font.color.rgb = COLOR_TEXT

        # 푸터
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        footer_p = footer_box.text_frame.paragraphs[0]
        footer_p.text = "Samsung Electronics HW Engineering Team | USB 3.2 Rev 1.1 Technical Seminar"
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = COLOR_DIM
        footer_p.alignment = PP_ALIGN.RIGHT

    # --- 슬라이드 생성 시작 ---

    # 1. 타이틀
    add_premium_title_slide("USB 3.2 Specification\nRevision 1.1", "Samsung Electronics HW Development Seminar\nAdvanced High-Speed Interface Deep Dive")

    # 2. Evolution
    add_premium_content_slide(1, "Architecture & Evolution", 
        ["USB 3.0/3.1을 USB 3.2 규격으로 통합 관리", 
         "Gen 1x2(10Gbps), Gen 2x2(20Gbps) 등 Multi-lane 기술 공식화", 
         "기존 PHY 기술 자산을 유지하며 대역폭 2배 확보 전략"],
        "USB 3.2 Rev 1.1은 기존 SuperSpeed(5G) 및 SuperSpeed+(10G) 설계를 기반으로 하며, Type-C의 가용 고속 라인을 병렬화하여 구현 시간을 단축함.")

    # 3. Multi-Lane
    add_premium_content_slide(2, "Multi-Lane Architecture (x2)", 
        ["Type-C 커넥터의 TX1/RX1, TX2/RX2 차동 쌍 동시 활용", 
         "Lane Bonding: 바이트 단위 Data Striping 및 분산 전송", 
         "수신단에서의 심볼 얼라이먼트(Symbol Alignment) 기술"],
        "Lane 0과 Lane 1로 나뉘어 전송되는 데이터는 수신단의 LTSSM 제어를 통해 완벽히 동기화되어야 하며, 이를 위해 전용 트레이닝 시퀀스가 사용됨.")

    # 4. Encoding
    add_premium_content_slide(3, "128b/132b Encoding Deep Dive", 
        ["Gen 2 규격 이상에서 적용되는 97% 고효율 인코딩", 
         "4-bit Sync Header를 통한 프레임 경계 및 클럭 복구", 
         "EMI 최소화를 위한 Scrambling LFSR 알고리즘"],
        "Sync Header (1010b=Data, 1100b=Control)는 8b/10b 인코딩의 20% 오버헤드를 3% 수준으로 획기적으로 개선하여 실효 전송률을 높임.")

    # 5. CC Logic
    add_premium_content_slide(4, "Type-C CC Interface Control", 
        ["Rd/Rp 전압 감지를 통한 삽입 및 방향(Flip) 판단", 
         "E-marker 케이블 식별 및 VCONN 전원 공급 자동화", 
         "PD VDM 메시지를 통한 x2 Multi-lane 진입 협상"],
        "고속 통신 중 PD 메시지 간섭을 방지하기 위해 CC 라인의 SI(Signal Integrity) 확보가 필수적이며, 특히 기생 커패시턴스 관리가 중요함.")

    # 6. De-skew
    add_premium_content_slide(5, "Lane-to-Lane De-skew Budget", 
        ["최대 허용 Skew: 6,400 ps (6.4ns) 규격 준수 필수", 
         "PHY 내부 Elastic Buffer를 이용한 동적 시간차 보정", 
         "PCB 패턴 및 케이블 지연 차이에 대한 마진 확보"],
        "삼성 HW 설계 가이드: Lane 간 Length Matching 편차를 0.25mm 이내로 유지하고, 고주파 손실이 적은 Low-loss PCB 소재 선정을 권장함.")

    # 7. LTSSM
    add_premium_content_slide(6, "LTSSM Link Training Steps", 
        ["Polling.Active: 비트 및 심볼 락 획득 과정", 
         "Polling.Configuration: Lane 수 및 속도 최종 협상", 
         "Rx.EQ: 전송 채널 특성 보정을 위한 3단계 등화"],
        "Recovery 상태는 링크 불안정 시 자동으로 진입하며, 재훈련 실패 시 하위 속도로의 Fallback을 통해 연결 안정성을 최우선으로 보장함.")

    # 8. Power Management
    add_premium_content_slide(7, "Power Management States", 
        ["U1/U2: 전력 절감과 빠른 복구(Latency)의 균형", 
         "U3: 최저 전력 Suspend 상태 (소프트웨어 제어)", 
         "LFPS: 저전력 상태 복귀를 위한 웨이크업 시그널링"],
        "LFPS 신호의 진폭 및 듀티 사이클이 규격을 벗어날 경우 절전 모드에서 시스템이 깨어나지 못하는 'Hang' 현상의 주요 원인이 됨.")

    # 9. SI/PI Guide
    add_premium_content_slide(8, "HW Signal Integrity Guidelines", 
        ["Differential Impedance: 90Ω ~ 100Ω 엄격 매칭", 
         "Insertion Loss: 5GHz(Nyquist) 기준 6dB 미만 권장", 
         "Crosstalk 차단을 위한 TX-RX 간격 최적화"],
        "Via Stub 제거를 위한 Back-drilling 기법 및 신호 전환부의 GND Stitching Via 배치는 20Gbps 설계의 필수 검토 항목임.")

    # 10. Conclusion
    add_premium_content_slide(9, "Final Implementation Summary", 
        ["20Gbps 인증 부품(PHY, Controller, Connector) 선정", 
         "Eye Diagram & Jitter (Rj, Dj) 분석을 통한 SI 검증", 
         "성공적인 삼성전자 차세대 제품 설계를 지원합니다."],
        "본 자료의 상세 설계 수치는 USB 3.2 Rev 1.1 표준 사양을 근거로 하며, 실제 보드 환경에 따른 SI 시뮬레이션 결과가 우선함.")

    # 저장
    filename = "USB_3_2_Premium_Technical_Seminar.pptx"
    prs.save(filename)
    print(f"성공: 고품격 디자인의 '{filename}' 파일이 생성되었습니다.")

if __name__ == "__main__":
    create_premium_presentation()
